#!/usr/bin/env python3
"""生成 SemAlign 项目答辩 PPT。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 配色：民航/企业科技风
COLOR_PRIMARY = RGBColor(0x00, 0x4B, 0x87)  # 深蓝
COLOR_ACCENT = RGBColor(0x00, 0x7A, 0xCC)  # 亮蓝
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SUB = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.2))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
    p2.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, section: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_ACCENT)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = section
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    sub_bullets: list[list[str]] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题栏
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_WHITE

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    sub_bullets = sub_bullets or [[] for _ in bullets]
    for i, (bullet, subs) in enumerate(zip(bullets, sub_bullets)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)
        for sub in subs:
            sp = tf.add_paragraph()
            sp.text = sub
            sp.level = 1
            sp.font.size = Pt(16)
            sp.font.color.rgb = COLOR_SUB
            sp.space_after = Pt(4)


def _add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_WHITE

    for col, (ctitle, items, x) in enumerate(
        [(left_title, left_items, 0.5), (right_title, right_items, 6.8)]
    ):
        h = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(5.8), Inches(0.5))
        hp = h.text_frame.paragraphs[0]
        hp.text = ctitle
        hp.font.size = Pt(22)
        hp.font.bold = True
        hp.font.color.rgb = COLOR_ACCENT
        box = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(5.8), Inches(5.2))
        tf = box.text_frame
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(17)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(6)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "民航多源标准治理系统",
        "SemAlign · 中航信标准治理系统\n基于大模型的标准语义理解、对齐与高效融合",
    )

    _add_content_slide(
        prs,
        "汇报提纲",
        [
            "建设背景与痛点",
            "系统定位与建设目标",
            "总体架构与技术路线",
            "核心功能演示路径",
            "对齐与冲突识别引擎",
            "质量保障与部署运维",
            "实现状态与演进规划",
        ],
    )

    _add_section_slide(prs, "01  建设背景")

    _add_content_slide(
        prs,
        "行业背景与痛点",
        [
            "标准来源多样：PDF、Word、Excel 及历史系统并存，检索困难",
            "表述不一致：同一概念跨标准命名差异大，人工理解成本高",
            "对齐工作量大：条款级差异、冲突与优先级难以系统化呈现",
            "治理缺乏闭环：导入、审核、发布、追溯缺少统一平台",
        ],
        sub_bullets=[
            ["分散在个人目录、邮件附件，难以形成企业资产池"],
            ["业务人员需反复翻阅原文比对"],
            ["两套标准并行时，人工逐条比对效率低"],
            ["标准生命周期与对齐任务状态难以追踪"],
        ],
    )

    _add_content_slide(
        prs,
        "课题目标",
        [
            "语义理解：文本向量化、混合检索、对齐助手对话",
            "标准对齐：条款切分、相似度匹配、冲突识别与优先级判定",
            "高效融合：人工审核发布、映射建议、工作台治理视图",
        ],
        sub_bullets=[
            ["gte-multilingual-base + BM25 + Chroma 混合召回"],
            ["四层冲突识别流水线 + 规则引擎"],
            ["draft → submitted → approved → published 状态机"],
        ],
    )

    _add_section_slide(prs, "02  系统定位")

    _add_content_slide(
        prs,
        "系统定位",
        [
            "企业内部标准治理与对齐融合平台（非官方发布系统）",
            "面向标准管理员、业务分析人员、系统管理员",
            "覆盖标准全生命周期：采集 → 解析 → 检索 → 对齐 → 审核发布",
        ],
    )

    _add_two_column_slide(
        prs,
        "用户角色与权限",
        "普通用户",
        [
            "工作台（只读指标）",
            "标准库浏览与下载",
            "智能检索与问答",
            "查看已发布对齐结果",
        ],
        "管理员",
        [
            "标准导入与批量入库",
            "发起对齐任务与冲突处理",
            "审核发布对齐成果",
            "用户与权限管理",
        ],
    )

    _add_section_slide(prs, "03  总体架构")

    _add_content_slide(
        prs,
        "逻辑架构（四层）",
        [
            "展示层：React SPA — 工作台 / 标准库 / 检索 / 导入 / 对齐",
            "应用服务层：FastAPI — REST API 与业务编排",
            "智能能力层：文档解析、向量检索、条款对齐、LLM 网关",
            "数据资源层：SQLite/PostgreSQL + 文件存储 + Chroma + BM25",
        ],
    )

    _add_two_column_slide(
        prs,
        "技术栈",
        "前端 semAlign",
        [
            "React 19 + TypeScript",
            "Vite 6 + React Router 7",
            "Tailwind CSS 4 + Zustand",
            "Recharts 数据可视化",
        ],
        "后端 semAlign_backend",
        [
            "FastAPI + SQLAlchemy",
            "Chroma + LangChain 混合检索",
            "gte-multilingual-base 嵌入",
            "OpenAI 兼容 LLM 网关（DeepSeek 等）",
        ],
    )

    _add_content_slide(
        prs,
        "核心业务流程",
        [
            "标准入库：上传 → 解析/OCR → 元数据抽取 → 校验 → 入库 → 向量索引",
            "智能检索：关键词 + 元数据 SQL + 向量语义 + BM25 混合召回",
            "标准对齐：选两组标准 → 条款切分 → 配对 → 冲突识别 → 人工决策",
            "审核发布：提交审核 → 批准/驳回 → 发布 → 普通用户可见",
        ],
    )

    _add_section_slide(prs, "04  核心功能")

    _add_content_slide(
        prs,
        "功能模块一览",
        [
            "标准导入：PDF/Excel 上传、解析预览、重复校验、导入历史",
            "标准数据库：列表筛选、详情查看、原文件/解析文本下载",
            "智能检索：元数据 + 全文 + 向量语义，可选 RAG 问答",
            "标准对齐：任务管理、冲突列表、人工 accept/reject/modify",
            "治理工作台：标准总量、月度新增、对齐任务、分类与动态",
            "对齐助手：基于 LLM 的条款解释与策略问答",
        ],
    )

    _add_content_slide(
        prs,
        "标准入库流程",
        [
            "POST /api/import/upload — 管理员上传 PDF/Excel（≤20MB）",
            "pdfplumber 文本层优先，乱码时 RapidOCR 回退",
            "抽取 GB/T、ISO 等标准号与元数据，校验重复/无效",
            "写入 standards 表，异步更新 Chroma 向量索引",
        ],
    )

    _add_content_slide(
        prs,
        "智能检索引擎",
        [
            "多通路召回：元数据模糊匹配 + Chroma 稠密向量 + BM25 稀疏检索",
            "嵌入模型：gte-multilingual-base，支持中英文混合标准文本",
            "EnsembleRetriever 融合排序，返回相关度与引用来源",
            "可选 RAG：结合检索片段调用 LLM 生成可解释回答",
        ],
    )

    _add_section_slide(prs, "05  对齐与冲突识别")

    _add_content_slide(
        prs,
        "四层冲突识别流水线",
        [
            "第一层 · 召回：候选条款对召回（相似度阈值过滤）",
            "第二层 · 智能体分析：原子差异、冲突置信度、权威性评分",
            "第三层 · 主动探寻：废止/替代关系、引用链 enrichment",
            "第四层 · 裁决输出：优先级排序、自然语言检测报告、解决方案",
        ],
        sub_bullets=[
            ["ClauseAligner / fallback 相似度矩阵"],
            ["layer2_agents 多维度分析"],
            ["layer3_active_seek 标准引用关系"],
            ["layer4_adjudication 融合为 comparison API 结构"],
        ],
    )

    _add_content_slide(
        prs,
        "优先级规则引擎",
        [
            "国际标准优先于国内标准",
            "最新修订版本优先",
            "强制性标准优先于推荐性标准",
            "综合性条款与语义置信度加权综合评分",
        ],
    )

    _add_content_slide(
        prs,
        "人工决策与审核发布",
        [
            "人工映射：accept / reject / modify + 备注回写",
            "审核状态机：draft → submitted → approved/rejected → published",
            "已发布对齐结果：普通用户只读访问",
            "用户反馈：赞成/反对、修改意见落库",
        ],
    )

    _add_section_slide(prs, "06  质量保障与部署")

    _add_content_slide(
        prs,
        "代码质量保障",
        [
            "SonarQube 9.9 LTS 静态代码分析，质量门禁 6/6 通过",
            "单元测试 124 项，新代码覆盖率 81.8%（≥80%）",
            "安全热点 100% 已审查，0 Bug / 0 漏洞",
            "pytest + coverage.xml 集成 CI 扫描流程",
        ],
    )

    _add_content_slide(
        prs,
        "部署架构",
        [
            "Docker Compose 一键部署：frontend + backend",
            "Nginx 反向代理：静态资源 + /api 转发",
            "数据持久化：./semAlign_backend/data 卷挂载",
            "访问地址：http://localhost:8080  |  API 文档 /docs",
        ],
    )

    _add_section_slide(prs, "07  实现状态与规划")

    _add_two_column_slide(
        prs,
        "实现状态对照",
        "已实现",
        [
            "前后端分离全栈",
            "PDF/Excel 导入解析",
            "混合检索与向量库",
            "四层对齐冲突流水线",
            "审核发布与权限体系",
            "对齐助手 LLM 对话",
            "Docker 生产部署",
        ],
        "规划中",
        [
            "LangGraph 多智能体编排",
            "Milvus 向量库演进",
            "Neo4j 知识图谱",
            "Scrapy 外部标准采集",
            "独立 Reranker 重排",
            "检索输入联想",
        ],
    )

    _add_content_slide(
        prs,
        "演进路线",
        [
            "近期：检索体验增强（Reranker、联想）、OCR 全书覆盖",
            "中期：知识图谱构建、术语关系可视化",
            "远期：LangGraph 编排、多源自动采集、与 OA/ERP 集成",
        ],
    )

    _add_content_slide(
        prs,
        "总结",
        [
            "SemAlign 形成「导入—检索—对齐—审核」标准治理闭环",
            "融合规则引擎 + 向量语义 + 大模型，兼顾准确性与可解释性",
            "工程化交付：Docker 部署、Sonar 质量门禁、完整设计文档",
            "为民航数据治理与跨标准合规分析提供可落地的知识底座",
        ],
    )

    # 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢聆听 · Q & A"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.8))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "SemAlign · 中航信标准治理系统  |  课题组"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

    return prs


def main() -> None:
    out_dir = Path(__file__).resolve().parents[1] / "semAlign" / "docs"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "SemAlign项目答辩PPT.pptx"
    prs = build_presentation()
    prs.save(str(out_path))
    print(f"已生成: {out_path}")
    print(f"幻灯片数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
